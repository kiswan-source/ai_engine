"""File readers: PDF, TXT, DOCX, CSV, JSON, Image."""
import errno, os, json, csv, time
from typing import Any, Callable, Dict, TypeVar

_T = TypeVar("_T")

# Fase 15 Decision 3 (failure recovery, Gate 1 Owner decision 2026-08-02):
# deliberately narrow — most read failures (corrupt file, wrong path,
# unsupported format) are deterministic, and retrying the identical
# operation wastes a round for nothing (a real Gate 1 finding: retry can't
# be driven by inspecting an already-stringified error message downstream,
# it has to see the real exception type where it's raised, i.e. here). Only
# genuinely transient OS-level conditions are retried; everything else either
# gets one format-specific fallback attempt (read_pdf/read_docx below) or
# fails immediately.
#
# Gate 2 finding (2026-08-02): the original exception-type-based version of
# this (BlockingIOError/InterruptedError/TimeoutError) was decorative —
# InterruptedError (EINTR) is auto-retried inside CPython's own I/O layer
# since PEP 475 and can never reach caller code; BlockingIOError only fires
# for non-blocking sockets, not local/mounted-drive file I/O; a real WSL
# 9p/drvfs hiccup surfaces as a plain OSError with EBUSY/EIO/EAGAIN/
# ETIMEDOUT. Fixed to check the real errno instead of the exception class —
# PermissionError/FileNotFoundError (also OSError subclasses, but
# deterministic — EACCES/ENOENT are deliberately NOT in this set) still fail
# on the first attempt, unchanged.
_TRANSIENT_ERRNOS = {errno.EAGAIN, errno.EBUSY, errno.EIO, errno.ETIMEDOUT}
if hasattr(errno, "EWOULDBLOCK"):
    _TRANSIENT_ERRNOS.add(errno.EWOULDBLOCK)
_RETRY_ATTEMPTS = 2
_RETRY_DELAY_SECONDS = 0.2


def _is_transient_os_error(exc: BaseException) -> bool:
    return isinstance(exc, OSError) and exc.errno in _TRANSIENT_ERRNOS


def _with_transient_retry(fn: Callable[[], _T]) -> _T:
    """Call ``fn()``, retrying up to ``_RETRY_ATTEMPTS`` more times only if
    it raises an ``OSError`` whose ``errno`` is genuinely transient (see
    ``_TRANSIENT_ERRNOS``). Any other exception — including a deterministic
    ``OSError`` like ``PermissionError``/``FileNotFoundError`` — propagates
    immediately on the first attempt, no wasted retries."""
    for attempt in range(_RETRY_ATTEMPTS + 1):
        try:
            return fn()
        except OSError as e:
            if not _is_transient_os_error(e) or attempt >= _RETRY_ATTEMPTS:
                raise
            time.sleep(_RETRY_DELAY_SECONDS * (attempt + 1))

# Fase 15 (DCF v5 mandate — bulk/large-file capability, Gate 1 Owner decision
# 2026-08-02): the 6 text-bearing readers below used to hard-truncate at
# 10,000 chars with no way to see the rest of a large document — the model
# just silently lost everything past that point. Pagination (offset/length,
# defaulting to the exact old 0..10000 window so every existing caller's
# default-args behavior is unchanged byte-for-byte) replaces that. Honest
# disclosure (Gate 1 finding, not fixed this pass — out of scope): this is a
# COMPLETENESS fix, not a performance one — read_docx/read_pptx/read_json
# still fully parse the whole file on every call regardless of which window
# is requested (python-docx/python-pptx/json.load have no random-access
# API); only the underlying full-text computation stays exactly as
# expensive as it always was, the change is only which slice of it gets
# returned. A real per-call cost reduction for those 3 formats would need a
# different library or an on-disk cache, deferred.
DEFAULT_PAGE_CHARS = 10000


def _paginate(full_text: str, offset: int, length: int) -> Dict[str, Any]:
    """Common offset/length windowing + pagination metadata for every
    text-bearing reader. ``offset``/``length`` default to 0/10000, the exact
    old hard-truncation window, so a caller that never passes them sees
    byte-identical `text`/`truncated`/`char_count` behavior to before this
    Fase."""
    offset = max(0, offset)
    length = max(0, length)
    window = full_text[offset : offset + length]
    has_more = offset + len(window) < len(full_text)
    return {
        "text": window,
        "char_count": len(full_text),
        "offset": offset,
        "returned_chars": len(window),
        "has_more": has_more,
        "truncated": has_more,  # kept for backward compatibility with pre-Fase-15 callers
    }


def _read_pdf_text(file_path: str, strict: bool) -> tuple[str, int]:
    from pypdf import PdfReader
    reader = PdfReader(file_path, strict=strict)
    pages = [(page.extract_text() or "").strip() for page in reader.pages]
    return "\n\n".join(t for t in pages if t), len(reader.pages)


def read_pdf(file_path: str, offset: int = 0, length: int = DEFAULT_PAGE_CHARS) -> Dict[str, Any]:
    if not os.path.exists(file_path):
        return {"error": f"File not found: {file_path}"}
    fallback_used = False
    try:
        full_text, page_count = _with_transient_retry(lambda: _read_pdf_text(file_path, strict=True))
    except ImportError:
        return {"error": "pypdf not installed"}
    except Exception as e:
        if _is_transient_os_error(e):
            return {"error": f"Gagal membaca {os.path.basename(file_path)} setelah beberapa percobaan "
                              f"(masalah sementara pada file/disk): {e}", "source": file_path}
        # Fallback (Gate 1 Decision 3 Option A): some malformed/oddly-encoded
        # PDFs fail pypdf's default strict parse but succeed leniently — a
        # real, cheap "try again more tolerantly" strategy, not a blind retry
        # of an identical call.
        try:
            full_text, page_count = _read_pdf_text(file_path, strict=False)
            fallback_used = True
        except Exception as e2:
            return {"error": str(e2), "source": file_path}
    result = {"source": file_path, "type": "pdf", "pages": page_count,
              **_paginate(full_text, offset, length)}
    if fallback_used:
        result["fallback_used"] = "lenient_parse"
    return result

def read_txt(file_path: str, offset: int = 0, length: int = DEFAULT_PAGE_CHARS) -> Dict[str, Any]:
    if not os.path.exists(file_path):
        return {"error": f"File not found: {file_path}"}
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        return {"source": file_path, "type": "txt", "line_count": content.count("\n"),
                **_paginate(content, offset, length)}
    except Exception as e:
        return {"error": str(e)}

def _read_docx_text(file_path: str) -> tuple[str, int]:
    from docx import Document
    doc = Document(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs), len(paragraphs)


def _read_docx_fallback_raw_xml(file_path: str):
    """Cheap fallback when python-docx's stricter parse fails (Gate 1
    Decision 3 Option A) — a .docx is a zip of XML parts; pull <w:t> text
    straight out of word/document.xml even when the structured python-docx
    parse can't open the file at all. Returns None (not raised) when even
    this can't recover anything, so the caller can report a clean error
    instead of a second traceback.

    Gate 2 finding (2026-08-02): the first version of this joined every
    <w:t> RUN with "\\n" — a real Word paragraph is routinely split across
    several runs (spellcheck, formatting boundaries), so a single sentence
    could come back with spurious line breaks inside it, silently
    (fallback_used was set, no error). Fixed to split on <w:p> PARAGRAPH
    boundaries first, joining runs within a paragraph with no separator and
    only paragraphs with "\\n" — matching _read_docx_text's own paragraph
    granularity exactly."""
    import re
    import zipfile
    import html

    try:
        with zipfile.ZipFile(file_path) as zf:
            xml_bytes = zf.read("word/document.xml")
    except Exception:
        return None
    xml_text = xml_bytes.decode("utf-8", errors="ignore")
    paragraph_blocks = re.findall(r"<w:p[ >].*?</w:p>", xml_text, flags=re.DOTALL)
    if not paragraph_blocks:
        return None
    lines = []
    for block in paragraph_blocks:
        runs = re.findall(r"<w:t[^>]*>(.*?)</w:t>", block, flags=re.DOTALL)
        line = "".join(html.unescape(r) for r in runs).strip()
        if line:
            lines.append(line)
    if not lines:
        return None
    return "\n".join(lines)


def read_docx(file_path: str, offset: int = 0, length: int = DEFAULT_PAGE_CHARS) -> Dict[str, Any]:
    if not os.path.exists(file_path):
        return {"error": f"File not found: {file_path}"}
    fallback_used = False
    paragraph_count = None
    try:
        full_text, paragraph_count = _with_transient_retry(lambda: _read_docx_text(file_path))
    except Exception as e:
        if _is_transient_os_error(e):
            return {"error": f"Gagal membaca {os.path.basename(file_path)} setelah beberapa percobaan "
                              f"(masalah sementara pada file/disk): {e}"}
        raw = _read_docx_fallback_raw_xml(file_path)
        if raw is None:
            return {"error": f"File .docx ini rusak atau bukan format Word yang valid, "
                              f"dan cara baca alternatif juga gagal."}
        full_text, fallback_used = raw, True
    result = {"source": file_path, "type": "docx", "paragraph_count": paragraph_count,
              **_paginate(full_text, offset, length)}
    if fallback_used:
        result["fallback_used"] = "raw_xml_extraction"
    return result

def read_csv(file_path: str) -> Dict[str, Any]:
    if not os.path.exists(file_path):
        return {"error": f"File not found: {file_path}"}
    try:
        rows = []
        with open(file_path, "r", encoding="utf-8", errors="ignore", newline="") as f:
            reader = csv.DictReader(f)
            headers = list(reader.fieldnames or [])
            for i, row in enumerate(reader):
                if i >= 500: break
                rows.append(dict(row))
        summary = ",".join(headers) + "\n" + "\n".join(",".join(str(v) for v in r.values()) for r in rows[:20])
        return {"source": file_path, "type": "csv", "headers": headers,
                "row_count": len(rows), "rows": rows[:50], "text": summary}
    except Exception as e:
        return {"error": str(e)}

def read_json(file_path: str, offset: int = 0, length: int = DEFAULT_PAGE_CHARS) -> Dict[str, Any]:
    if not os.path.exists(file_path):
        return {"error": f"File not found: {file_path}"}
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        text = json.dumps(data, indent=2, ensure_ascii=False)
        result = {"source": file_path, "type": "json", **_paginate(text, offset, length)}
        if offset == 0:
            # Backward compatible: the parsed object was always returned in
            # full pre-Fase-15 (only `text` was windowed) — keep doing that
            # for the first page only, since re-sending the whole parsed
            # `data` on every subsequent page would defeat the point of
            # pagination.
            result["data"] = data
        return result
    except Exception as e:
        return {"error": str(e)}

def read_xlsx(file_path: str, offset: int = 0, length: int = DEFAULT_PAGE_CHARS) -> Dict[str, Any]:
    """Workspace Slice 3 (Fase 12) — per-sheet row preview, same shape as
    read_csv's summary but keyed per sheet since a workbook can have more
    than one."""
    if not os.path.exists(file_path):
        return {"error": f"File not found: {file_path}"}
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        sheet_names = list(wb.sheetnames)
        row_counts: Dict[str, int] = {}
        text_parts = []
        for sheet_name in sheet_names:
            ws = wb[sheet_name]
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= 500:
                    break
                rows.append(row)
            row_counts[sheet_name] = len(rows)
            preview = "\n".join(",".join("" if v is None else str(v) for v in r) for r in rows[:20])
            text_parts.append(f"[Sheet: {sheet_name}] ({len(rows)} baris)\n{preview}")
        wb.close()
        full_text = "\n\n".join(text_parts)
        return {"source": file_path, "type": "xlsx", "sheet_names": sheet_names,
                "sheet_count": len(sheet_names), "row_counts": row_counts,
                **_paginate(full_text, offset, length)}
    except ImportError:
        return {"error": "openpyxl not installed"}
    except Exception as e:
        return {"error": str(e), "source": file_path}


def read_pptx(file_path: str, offset: int = 0, length: int = DEFAULT_PAGE_CHARS) -> Dict[str, Any]:
    """Workspace Slice 3 (Fase 12) — per-slide text, same paragraph-joining
    approach read_docx already uses for a Word document's paragraphs."""
    if not os.path.exists(file_path):
        return {"error": f"File not found: {file_path}"}
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides):
            lines = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs).strip()
                        if text:
                            lines.append(text)
            slides_text.append(f"[Slide {i + 1}]\n" + "\n".join(lines))
        full_text = "\n\n".join(slides_text)
        return {"source": file_path, "type": "pptx", "slide_count": len(prs.slides),
                **_paginate(full_text, offset, length)}
    except ImportError:
        return {"error": "python-pptx not installed"}
    except Exception as e:
        return {"error": str(e), "source": file_path}


def read_image(file_path: str) -> Dict[str, Any]:
    if not os.path.exists(file_path):
        return {"error": f"File not found: {file_path}"}
    try:
        from PIL import Image
        img = Image.open(file_path)
        w, h = img.size
        desc = f"Image: {w}x{h}px, mode={img.mode}"
        try:
            import pytesseract
            ocr = pytesseract.image_to_string(img, lang="eng+ind").strip()
        except Exception:
            ocr = ""
        return {"source": file_path, "type": "image", "description": desc,
                "width": w, "height": h,
                "text": ocr if ocr else f"[Image: {os.path.basename(file_path)}, no OCR text]",
                "has_ocr_text": bool(ocr)}
    except ImportError:
        return {"error": "Pillow not installed"}
    except Exception as e:
        return {"error": str(e)}


# Fase 15 (bulk-file capability, Gate 1 Owner decision 2026-08-02): reading
# N files used to need N separate model-directed tool calls. One extension
# -> reader mapping, mirroring the same mapping ToolRegistry.auto_reader()
# already uses per-extension elsewhere in this codebase, so file-type
# dispatch logic isn't duplicated with a different set of rules.
_EXTENSION_READERS = {
    "pdf": read_pdf,
    "txt": read_txt, "md": read_txt, "log": read_txt,
    "docx": read_docx, "doc": read_docx,
    "csv": read_csv,
    "json": read_json,
    "xlsx": read_xlsx,
    "pptx": read_pptx,
    "jpg": read_image, "jpeg": read_image, "png": read_image, "webp": read_image,
    "tif": read_image, "tiff": read_image, "bmp": read_image, "gif": read_image,
}
# The 6 readers that accept offset/length (agent/tools/readers.py's own
# pagination, above) — read_csv/read_image have their own unrelated shapes.
_PAGINATED_READERS = {read_pdf, read_txt, read_docx, read_json, read_xlsx, read_pptx}

# Gate 1 finding: an unbounded file_paths list is a real DoS vector (memory/
# time) — capped, not a soft suggestion. Each individual read still costs
# whatever that one file's full-parse already costs (Decision 2's "not a
# performance fix" caveat applies per-file here too) — this cap bounds the
# NUMBER of files, not the cost of any single large one.
MAX_BATCH_FILES = 50
# Gate 2 finding (2026-08-02): length_per_file had no upper clamp — a caller
# passing a huge value defeated the whole point of a small per-file budget
# in a batch call (response size, and for read_pdf/read_xlsx more of the
# already-computed full text gets serialized). Clamped, not just defaulted.
MAX_LENGTH_PER_FILE = 5000


def read_many_files(file_paths, length_per_file: int = 2000) -> Dict[str, Any]:
    """Read several files in one call instead of one model-directed tool
    call per file. Dispatches each path to the matching single-file reader
    by extension; a per-file result carries an "error" key on failure (same
    shape every single-file reader already uses) rather than aborting the
    whole batch — one bad file must not lose every other file's result."""
    if isinstance(file_paths, str):
        file_paths = [file_paths]
    if not file_paths:
        return {"success": False, "error": "file_paths kosong."}
    if len(file_paths) > MAX_BATCH_FILES:
        return {
            "success": False,
            "error": f"Terlalu banyak file sekaligus ({len(file_paths)}), maksimum "
                     f"{MAX_BATCH_FILES} per panggilan. Bagi jadi beberapa panggilan.",
        }
    length_per_file = max(1, min(length_per_file, MAX_LENGTH_PER_FILE))

    results = []
    for file_path in file_paths:
        ext = os.path.splitext(file_path)[1].lstrip(".").lower()
        reader = _EXTENSION_READERS.get(ext)
        if reader is None:
            results.append({"file_path": file_path, "error": f"Ekstensi '.{ext}' tidak didukung untuk baca batch."})
            continue
        try:
            if reader in _PAGINATED_READERS:
                item = reader(file_path, offset=0, length=length_per_file)
            else:
                item = reader(file_path)
        except Exception as e:  # a single bad file must not abort the whole batch
            item = {"error": str(e)}
        results.append({"file_path": file_path, **item})

    ok_count = sum(1 for r in results if "error" not in r)
    return {
        "success": True, "count": len(results), "ok_count": ok_count,
        "text": f"{ok_count}/{len(results)} file berhasil dibaca.",
        "results": results,
    }
