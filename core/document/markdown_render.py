"""Shared Markdown -> structured-blocks parser for document generation
(Fase 11 — "format profesional" fix).

`agent/tools/writers.py::write_docx`/`write_pdf` used to hand-parse content
line by line (`"# "`/`"## "` prefix match, `line.split("**")` for bold) —
that missed numbered lists, tables, nested formatting inside a heading, and
italic/inline-code entirely, so a heading like ``### **Ringkasan**`` or any
GFM table the model produced (the system prompt explicitly asks for both)
came out with literal ``###``/``**``/``|`` characters still in the generated
document. Same root cause the frontend chat renderer fix (Fase 10,
`web/src/components/chat/ChatMarkdown.tsx`) already closed for the chat
bubble — this is the Python-ecosystem equivalent, one parse (`markdown-it-py`,
CommonMark + GFM tables) feeding both `write_docx` (python-docx) and
`write_pdf` (ReportLab) so there's exactly one place that understands
Markdown, not two hand-rolled parsers that drift apart.

Deliberately a plain intermediate representation (`Block`/`Run` dataclasses)
rather than handing raw HTML to either library — python-docx has no HTML
importer at all, and ReportLab's `Paragraph` only understands a small
hand-picked HTML-like subset, so a real AST walk is the only approach that
works for both.
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from typing import Literal

from markdown_it import MarkdownIt

_MD = MarkdownIt("commonmark").enable("table")


@dataclass
class Run:
    text: str
    bold: bool = False
    italic: bool = False
    code: bool = False


BlockKind = Literal[
    "heading", "paragraph", "bullet_list", "ordered_list", "table", "hr", "blockquote", "code_block"
]


@dataclass
class Block:
    kind: BlockKind
    level: int = 0
    runs: list[Run] = field(default_factory=list)
    # (depth, runs) per item — depth is 1 for a top-level item, 2+ for
    # nested items (Gate 2 fix: nested items used to be dropped outright
    # rather than just flattened, see parse_markdown below).
    items: list[tuple[int, list[Run]]] = field(default_factory=list)
    header: list[list[Run]] = field(default_factory=list)
    rows: list[list[list[Run]]] = field(default_factory=list)
    text: str = ""
    start: int = 1


def _inline_runs(inline_token) -> list[Run]:
    """Flatten one `inline` token's children into a flat run list — nested
    bold/italic/code become independent boolean flags per run rather than a
    tree, since neither python-docx nor ReportLab needs anything richer than
    "this span is bold/italic/code" (no bold-within-italic-within-code
    distinction the model would plausibly produce in a report)."""
    runs: list[Run] = []
    bold = italic = False
    for c in inline_token.children or []:
        if c.type == "text":
            if c.content:
                runs.append(Run(c.content, bold, italic))
        elif c.type == "code_inline":
            runs.append(Run(c.content, bold, italic, code=True))
        elif c.type in ("softbreak", "hardbreak"):
            runs.append(Run("\n", bold, italic))
        elif c.type == "strong_open":
            bold = True
        elif c.type == "strong_close":
            bold = False
        elif c.type == "em_open":
            italic = True
        elif c.type == "em_close":
            italic = False
        elif c.type == "image":
            # Gate 2 fix: used to silently produce zero runs (fell through
            # every branch above) — no image-embedding support exists (would
            # need to fetch/decode arbitrary URLs into both python-docx and
            # ReportLab), but dropping the content with no trace is worse
            # than a visible placeholder. Same "plain text, not the real
            # thing" treatment already applied to links/strikethrough below.
            alt = c.content or (c.attrs.get("alt", "") if c.attrs else "")
            runs.append(Run(f"[Gambar: {alt}]" if alt else "[Gambar]", bold, italic))
        # link_open/close, s_open/close (strikethrough): rendered as plain
        # text runs — a generated report doesn't need clickable links, and
        # strikethrough has no natural equivalent worth the complexity here.
    return runs


def parse_markdown(content: str) -> list[Block]:
    """One markdown-it-py pass -> a flat list of `Block`s, in document order."""
    tokens = _MD.parse(content)
    blocks: list[Block] = []
    i, n = 0, len(tokens)
    while i < n:
        t = tokens[i]
        if t.type == "heading_open":
            blocks.append(Block(kind="heading", level=int(t.tag[1]), runs=_inline_runs(tokens[i + 1])))
            i += 3
        elif t.type == "paragraph_open":
            blocks.append(Block(kind="paragraph", runs=_inline_runs(tokens[i + 1])))
            i += 3
        elif t.type in ("bullet_list_open", "ordered_list_open"):
            kind: BlockKind = "bullet_list" if t.type == "bullet_list_open" else "ordered_list"
            start = 1
            if t.type == "ordered_list_open" and t.attrs and t.attrs.get("start") is not None:
                start = int(t.attrs["start"])
            # Gate 2 fix: nested items (depth > 1) used to be skipped
            # outright rather than just flattened — a genuine content-loss
            # bug, not a cosmetic one. Now every item is kept, tagged with
            # its depth so both renderers can indent it instead of losing it.
            items: list[tuple[int, list[Run]]] = []
            depth = 1
            i += 1
            while i < n and depth > 0:
                tok = tokens[i]
                if tok.type in ("bullet_list_open", "ordered_list_open"):
                    depth += 1
                elif tok.type in ("bullet_list_close", "ordered_list_close"):
                    depth -= 1
                    if depth == 0:
                        i += 1
                        break
                elif tok.type == "inline":
                    items.append((depth, _inline_runs(tok)))
                i += 1
            blocks.append(Block(kind=kind, items=items, start=start))
        elif t.type == "table_open":
            header: list[list[Run]] = []
            rows: list[list[list[Run]]] = []
            current_row: list[list[Run]] | None = None
            in_body = False
            i += 1
            while i < n and tokens[i].type != "table_close":
                tok = tokens[i]
                if tok.type == "tbody_open":
                    in_body = True
                elif tok.type == "tr_open":
                    current_row = []
                elif tok.type == "tr_close" and current_row is not None:
                    if in_body:
                        rows.append(current_row)
                    else:
                        header = current_row
                    current_row = None
                elif tok.type == "inline" and current_row is not None:
                    current_row.append(_inline_runs(tok))
                i += 1
            blocks.append(Block(kind="table", header=header, rows=rows))
            i += 1
        elif t.type == "blockquote_open":
            runs: list[Run] = []
            i += 1
            while i < n and tokens[i].type != "blockquote_close":
                if tokens[i].type == "inline":
                    if runs:
                        runs.append(Run("\n"))
                    runs.extend(_inline_runs(tokens[i]))
                i += 1
            blocks.append(Block(kind="blockquote", runs=runs))
            i += 1
        elif t.type in ("fence", "code_block"):
            blocks.append(Block(kind="code_block", text=t.content.rstrip("\n")))
            i += 1
        elif t.type == "hr":
            blocks.append(Block(kind="hr"))
            i += 1
        else:
            i += 1
    return blocks


def _add_docx_blocks(blocks: list[Block], add_paragraph, add_table) -> None:
    """Shared block-rendering logic for docx, parameterized over WHERE new
    paragraphs/tables land: ``add_paragraph(text="", style=None)`` and
    ``add_table(cols)`` abstract over "append at the end of the document"
    (render_docx_body) vs "insert at a specific position" (edit_docx_section,
    Workspace Slice 4, Fase 12) — one block-kind dispatch, not two copies
    drifting apart. Signature-compatible with both `Document.add_paragraph`
    and `Paragraph.insert_paragraph_before`, so either can be passed directly
    with no wrapper needed for the "append" case."""
    from docx.shared import Pt

    def add_runs(paragraph, runs: list[Run]) -> None:
        for r in runs:
            if r.text == "\n":
                paragraph.add_run().add_break()
                continue
            run = paragraph.add_run(r.text)
            run.bold = r.bold
            run.italic = r.italic
            if r.code:
                run.font.name = "Consolas"
                run.font.size = Pt(10)

    for block in blocks:
        if block.kind == "heading":
            level = min(max(block.level, 1), 4)
            add_runs(add_paragraph(style=f"Heading {level}"), block.runs)
        elif block.kind == "paragraph":
            add_runs(add_paragraph(), block.runs)
        elif block.kind == "bullet_list":
            for depth, item_runs in block.items:
                style = f"List Bullet {min(depth, 3)}" if depth > 1 else "List Bullet"
                add_runs(add_paragraph(style=style), item_runs)
        elif block.kind == "ordered_list":
            for depth, item_runs in block.items:
                style = f"List Number {min(depth, 3)}" if depth > 1 else "List Number"
                add_runs(add_paragraph(style=style), item_runs)
        elif block.kind == "table":
            n_cols = len(block.header) if block.header else (len(block.rows[0]) if block.rows else 0)
            if n_cols == 0:
                continue
            table = add_table(n_cols)
            table.style = "Light Grid Accent 1"
            if block.header:
                cells = table.add_row().cells
                for i, cell_runs in enumerate(block.header):
                    add_runs(cells[i].paragraphs[0], cell_runs)
                    for run in cells[i].paragraphs[0].runs:
                        run.bold = True
            for row in block.rows:
                cells = table.add_row().cells
                for i, cell_runs in enumerate(row):
                    if i < n_cols:
                        add_runs(cells[i].paragraphs[0], cell_runs)
        elif block.kind == "blockquote":
            p = add_paragraph()
            p.paragraph_format.left_indent = Pt(18)
            add_runs(p, block.runs)
        elif block.kind == "code_block":
            p = add_paragraph()
            run = p.add_run(block.text)
            run.font.name = "Consolas"
            run.font.size = Pt(9)
        elif block.kind == "hr":
            add_paragraph("─" * 40)


def render_docx_body(doc, content: str) -> None:
    """Append parsed ``content`` to an existing python-docx ``Document`` —
    called after the title/heading-0 block `write_docx` already added, so
    markdown headings map to docx levels 1-4 (never 0, reserved for the
    document title)."""
    _add_docx_blocks(parse_markdown(content), doc.add_paragraph, lambda cols: doc.add_table(rows=0, cols=cols))


def edit_docx_section(doc, heading_text: str, content: str, heading_level: int | None = None) -> str:
    """In-place section edit (Workspace Slice 4, Fase 12) — replaces the
    BODY of the section under the paragraph whose text matches
    ``heading_text`` (case-insensitive, a "Heading 1"-"Heading 4" style
    paragraph — the first one found in document order, or a specific level
    if ``heading_level`` disambiguates) with freshly-rendered ``content``.
    "Body" is everything from just after that heading up to (not including)
    the next heading at the same or higher level, or the end of the
    document if there is none — the heading paragraph itself is left
    untouched, only what's under it changes. If no paragraph matches
    ``heading_text`` (create-if-missing), the heading + content are
    appended at the end of the document instead. Returns "edited" or
    "appended".

    Owner decision (Gate 1, Fase 12): this is the DOCX half of Workspace
    Slice 4 — full mid-document PDF text editing was deliberately NOT built
    the same way (would need PyMuPDF, AGPL-3.0-licensed, classified
    HUMAN-ONLY/HALT by the DCF decision engine over the license exposure).
    PDF instead gets a narrower append-only operation — see
    `agent/tools/workspace_reader.py::_append_pdf_in_place`.
    """
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph

    # Gate 2 fix: the chat-facing wrapper (workspace_write_file) already
    # rejects a blank heading before this is ever called, but this function
    # is itself public — called directly (as it is by its own test suite),
    # a blank/whitespace-only heading_text would previously match nothing,
    # fall to the "not found" branch, and silently append a real Heading 1
    # paragraph with EMPTY text. Guard here too rather than relying solely
    # on one caller's validation.
    if not heading_text or not heading_text.strip():
        raise ValueError("heading_text tidak boleh kosong")

    def heading_level_of(paragraph) -> int | None:
        style_name = paragraph.style.name if paragraph.style else ""
        if not style_name.startswith("Heading"):
            return None
        try:
            return int(style_name.rsplit(" ", 1)[-1])
        except ValueError:
            return None

    needle = heading_text.strip().lower()
    target, target_level = None, None
    for p in doc.paragraphs:
        level = heading_level_of(p)
        if level is None or p.text.strip().lower() != needle:
            continue
        if heading_level is not None and level != heading_level:
            continue
        target, target_level = p, level
        break

    blocks = parse_markdown(content)

    if target is None:
        level = min(max(heading_level or 1, 1), 4)
        doc.add_paragraph(heading_text, style=f"Heading {level}")
        _add_docx_blocks(blocks, doc.add_paragraph, lambda cols: doc.add_table(rows=0, cols=cols))
        return "appended"

    # Delete every paragraph/table between the target heading and the next
    # same-or-higher-level heading — a deeper heading (e.g. an H2 under this
    # H1) is part of the section being replaced, not a boundary, so it's
    # removed along with everything else. Found by adversarial review: when
    # the target is the LAST section (no boundary heading found, end_idx
    # defaults to len(children)), deleting EVERY body child in range used to
    # include the document's trailing <w:sectPr> (page size/margins/
    # header-footer refs, always the final body child) — silently dropping
    # it corrupts the document (doc.sections goes from 1 to 0, survives
    # save+reload). Restricting the sweep to w:p/w:tbl only — the two
    # element types this section-replace ever creates or means to remove —
    # leaves sectPr (and any other structural element) untouched regardless
    # of whether a boundary was found.
    children = list(doc.element.body)
    start_idx = children.index(target._p)
    boundary_el = None
    for el in children[start_idx + 1:]:
        if el.tag != qn("w:p"):
            continue
        lvl = heading_level_of(Paragraph(el, doc))
        if lvl is not None and lvl <= target_level:
            boundary_el = el
            break
    end_idx = children.index(boundary_el) if boundary_el is not None else len(children)
    for el in children[start_idx + 1:end_idx]:
        if el.tag in (qn("w:p"), qn("w:tbl")):
            el.getparent().remove(el)

    if boundary_el is not None:
        ref_paragraph = Paragraph(boundary_el, doc)

        def add_paragraph(text="", style=None):
            return ref_paragraph.insert_paragraph_before(text, style)

        def add_table(cols):
            table = doc.add_table(rows=0, cols=cols)
            ref_paragraph._p.addprevious(table._tbl)
            return table
    else:
        # Target was the last section in the document — inserting "before
        # the next heading" has no anchor, so appending at the very end of
        # the document body is equivalent (nothing comes after it anyway).
        add_paragraph = doc.add_paragraph
        add_table = lambda cols: doc.add_table(rows=0, cols=cols)  # noqa: E731

    _add_docx_blocks(blocks, add_paragraph, add_table)
    return "edited"


def esc_pdf_markup(text: str) -> str:
    """Escape plain text for ReportLab's Paragraph mini-XML markup — module-
    level (not a render_pdf_story-local closure) so any caller that hands
    literal text to a ReportLab Paragraph outside a parsed Block (e.g.
    write_pdf's/append_pdf_section's `title` argument, Gate 2 fix) can
    reuse it instead of duplicating this or, worse, skipping it. An
    unescaped `&`/`<`/`>` either breaks ReportLab's markup parser outright
    (a real inline tag left unclosed) or silently drops the offending
    substring (parsed as an unknown empty tag)."""
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def render_pdf_story(content: str, styles: dict) -> list:
    """Parsed ``content`` -> a list of ReportLab flowables. ``styles`` must
    provide ``body``/``h1``/``h2``/``bullet``/``quote``/``code``/
    ``table_header``/``table_cell`` `ParagraphStyle`s — built by
    `write_pdf` so the accent color/fonts stay consistent with its title."""
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import HRFlowable, Paragraph, Spacer, Table, TableStyle

    esc = esc_pdf_markup

    def heading_style(level: int):
        if level <= 1:
            return styles["h1"]
        if level == 2:
            return styles["h2"]
        # Gate 2 fix: `styles` only carries h1/h2 (write_pdf's fixed set) —
        # used to collapse every level >= 2 into h2, so H2/H3/H4 all looked
        # identical despite docx keeping them visually distinct. Derive a
        # progressively smaller variant instead of adding new required style
        # keys to every write_pdf caller.
        base = styles["h2"]
        shrink = 2 * (level - 2)
        return ParagraphStyle(
            f"h{level}_derived", parent=base,
            fontSize=max(base.fontSize - shrink, 9), leading=max(base.leading - shrink, 11),
        )

    def indented_style(depth: int):
        if depth <= 1:
            return styles["bullet"]
        return styles["bullet"].clone(
            f"bullet_d{depth}", leftIndent=(styles["bullet"].leftIndent or 0) + (depth - 1) * 0.6 * cm
        )

    def runs_to_markup(runs: list[Run]) -> str:
        parts = []
        for r in runs:
            if r.text == "\n":
                parts.append("<br/>")
                continue
            text = esc(r.text)
            if r.code:
                text = f'<font face="Courier">{text}</font>'
            if r.bold:
                text = f"<b>{text}</b>"
            if r.italic:
                text = f"<i>{text}</i>"
            parts.append(text)
        return "".join(parts)

    story: list = []
    for block in parse_markdown(content):
        if block.kind == "heading":
            story.append(Paragraph(runs_to_markup(block.runs), heading_style(block.level)))
        elif block.kind == "paragraph":
            story.append(Paragraph(runs_to_markup(block.runs), styles["body"]))
        elif block.kind in ("bullet_list", "ordered_list"):
            # Gate 2 fix: honor the list's actual start number (markdown
            # `3. foo`) and indent nested items instead of dropping them —
            # counters are per-depth so a nested numbered sub-list restarts
            # at 1 rather than continuing its parent's count.
            counters: dict[int, int] = {}
            for depth, item_runs in block.items:
                for d in [d for d in counters if d > depth]:
                    del counters[d]
                if block.kind == "ordered_list":
                    base = block.start if depth == 1 else 1
                    counters[depth] = counters.get(depth, base - 1) + 1
                    prefix = f"{counters[depth]}."
                else:
                    prefix = "•"
                story.append(Paragraph(f"{prefix} {runs_to_markup(item_runs)}", indented_style(depth)))
        elif block.kind == "table":
            data = []
            if block.header:
                data.append([Paragraph(runs_to_markup(c), styles["table_header"]) for c in block.header])
            for row in block.rows:
                data.append([Paragraph(runs_to_markup(c), styles["table_cell"]) for c in row])
            if data:
                tbl = Table(data, hAlign="LEFT")
                cmds = [
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#dee2e6")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ]
                if block.header:
                    cmds.append(("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#00c896")))
                    cmds.append(("TEXTCOLOR", (0, 0), (-1, 0), colors.white))
                tbl.setStyle(TableStyle(cmds))
                story.append(tbl)
        elif block.kind == "blockquote":
            story.append(Paragraph(runs_to_markup(block.runs), styles["quote"]))
        elif block.kind == "code_block":
            story.append(Paragraph(esc(block.text).replace("\n", "<br/>"), styles["code"]))
        elif block.kind == "hr":
            story.append(HRFlowable(width="100%", thickness=0.75, color=colors.HexColor("#dee2e6")))
        story.append(Spacer(1, 0.15 * cm))
    return story


def _runs_to_plain_text(runs: list[Run]) -> str:
    """Flatten runs to plain text — xlsx cells and pptx text frames don't
    need the bold/italic/code distinction render_docx_body/render_pdf_story
    preserve, just the actual characters."""
    return "".join("\n" if r.text == "\n" else r.text for r in runs)


def _group_by_top_heading(blocks: list[Block], default_title: str) -> list[tuple[str, list[Block]]]:
    """Split a parsed block list into (name, blocks) groups at every H1 —
    shared by render_xlsx_workbook (H1 -> new sheet) and render_pptx_slides
    (H1 -> new slide), both Workspace Slice 3 (Fase 12). Content before the
    first H1 (or the whole document, if there's no H1 at all) is grouped
    under ``default_title``; only that one leading group is dropped when
    empty (document starts right at an H1, so there's nothing to put under
    ``default_title``). Every H1 AFTER the first one always gets its own
    group even if empty — an outline heading with no body yet (a realistic
    shape of model output) still needs its own sheet/slide, not to vanish
    silently (found by adversarial review: the original version only kept
    non-empty groups, so any H1 immediately followed by another H1 dropped
    with zero trace, not just a genuinely leading empty one)."""
    groups: list[tuple[str, list[Block]]] = []
    current_name = default_title
    current_blocks: list[Block] = []
    seen_heading = False
    for block in blocks:
        if block.kind == "heading" and block.level <= 1:
            if current_blocks or seen_heading:
                groups.append((current_name, current_blocks))
            current_name = _runs_to_plain_text(block.runs) or "Sheet"
            current_blocks = []
            seen_heading = True
        else:
            current_blocks.append(block)
    groups.append((current_name, current_blocks))
    return groups


def render_xlsx_workbook(content: str, default_title: str = "Sheet1"):
    """Parsed ``content`` -> an `openpyxl.Workbook` (Workspace Slice 3,
    Fase 12) — each top-level (H1) heading starts a new sheet; content
    before the first one (or the whole document, if there's no H1) lands on
    a sheet named ``default_title``. GFM tables become real multi-column
    rows (the main reason this exists — a table is what actually belongs in
    a spreadsheet); paragraphs/lists/headings become single-cell rows,
    reasonable for the common "generate a data table with commentary" case
    but not a general-purpose docx-to-xlsx converter."""
    from openpyxl import Workbook
    from openpyxl.styles import Font

    def safe_sheet_name(name: str) -> str:
        cleaned = re.sub(r'[\\/*?:\[\]]', "-", name).strip() or "Sheet"
        return cleaned[:31]

    bold = Font(bold=True)
    wb = Workbook()
    wb.remove(wb.active)
    seen_names: set[str] = set()

    for name, group_blocks in _group_by_top_heading(parse_markdown(content), default_title):
        sheet_name = safe_sheet_name(name)
        base, n = sheet_name, 2
        while sheet_name in seen_names:
            # Compute the truncation from the actual suffix length rather
            # than a fixed budget — found by adversarial review: a fixed
            # `base[:28]` stays within Excel's 31-char limit only while `n`
            # is 1-2 digits; heavy duplication (n >= 100) would silently
            # produce an invalid, spec-violating sheet name otherwise.
            suffix = f"-{n}"
            sheet_name = f"{base[:31 - len(suffix)]}{suffix}"
            n += 1
        seen_names.add(sheet_name)
        ws = wb.create_sheet(sheet_name)

        for block in group_blocks:
            if block.kind == "heading":
                ws.append([_runs_to_plain_text(block.runs)])
                for c in ws[ws.max_row]:
                    c.font = bold
            elif block.kind == "paragraph":
                ws.append([_runs_to_plain_text(block.runs)])
            elif block.kind in ("bullet_list", "ordered_list"):
                for depth, item_runs in block.items:
                    ws.append(["  " * (depth - 1) + "- " + _runs_to_plain_text(item_runs)])
            elif block.kind == "table":
                if block.header:
                    ws.append([_runs_to_plain_text(c) for c in block.header])
                    for c in ws[ws.max_row]:
                        c.font = bold
                for row in block.rows:
                    ws.append([_runs_to_plain_text(c) for c in row])
            elif block.kind == "blockquote":
                ws.append([_runs_to_plain_text(block.runs)])
            elif block.kind == "code_block":
                ws.append([block.text])
            elif block.kind == "hr":
                ws.append([])

    if not wb.sheetnames:
        wb.create_sheet(safe_sheet_name(default_title))
    return wb


def render_pptx_slides(content: str, default_title: str = "Ringkasan"):
    """Parsed ``content`` -> a `pptx.Presentation` (Workspace Slice 3, Fase
    12) — each top-level (H1) heading starts a new slide (title + a bullet
    body from any headings/paragraphs/lists in that section); a GFM table
    gets its own separate "Title Only" slide with a real pptx table shape
    (mixing a table into a bulleted body placeholder isn't a sensible
    layout, so it's deliberately split out rather than forced in)."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    body_layout = prs.slide_layouts[1]  # "Title and Content"
    table_layout = prs.slide_layouts[5]  # "Title Only"

    for name, group_blocks in _group_by_top_heading(parse_markdown(content), default_title):
        slide = prs.slides.add_slide(body_layout)
        slide.shapes.title.text = name
        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()
        first_line = True
        table_blocks = []

        def add_line(text: str, level: int = 0) -> None:
            nonlocal first_line
            if first_line:
                text_frame.text = text
                para = text_frame.paragraphs[0]
                first_line = False
            else:
                para = text_frame.add_paragraph()
                para.text = text
            para.level = level

        for block in group_blocks:
            if block.kind == "table":
                table_blocks.append(block)
                continue
            if block.kind in ("bullet_list", "ordered_list"):
                # Gate 2 fix: list items used to render as bare text with
                # no marker at all (not even a literal "•"/"N.", unlike the
                # xlsx/PDF renderers built from the same IR) and depth was
                # marked only by literal leading spaces rather than
                # `paragraph.level` — so PowerPoint's own per-level bullet/
                # indent styling never applied either. Per-depth counters
                # (reset on returning to a shallower depth) match the same
                # pattern render_pdf_story already uses for ordered lists.
                counters: dict[int, int] = {}
                for depth, item_runs in block.items:
                    for d in [d for d in counters if d > depth]:
                        del counters[d]
                    if block.kind == "ordered_list":
                        counters[depth] = counters.get(depth, 0) + 1
                        prefix = f"{counters[depth]}. "
                    else:
                        prefix = "• "
                    add_line(prefix + _runs_to_plain_text(item_runs), level=min(depth - 1, 8))
                continue
            if block.kind in ("heading", "paragraph", "blockquote"):
                lines = [_runs_to_plain_text(block.runs)]
            elif block.kind == "code_block":
                lines = [block.text]
            else:  # hr — no textual content to add to the body
                lines = []
            for line in lines:
                add_line(line)

        for tbl_block in table_blocks:
            n_cols = len(tbl_block.header) if tbl_block.header else (len(tbl_block.rows[0]) if tbl_block.rows else 0)
            if n_cols == 0:
                continue
            n_rows = (1 if tbl_block.header else 0) + len(tbl_block.rows)
            table_slide = prs.slides.add_slide(table_layout)
            table_slide.shapes.title.text = f"{name} — Tabel"
            shape = table_slide.shapes.add_table(
                n_rows, n_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(min(0.6 * n_rows, 5.5))
            )
            table = shape.table
            r = 0
            if tbl_block.header:
                for c, cell_runs in enumerate(tbl_block.header):
                    table.cell(0, c).text = _runs_to_plain_text(cell_runs)
                r = 1
            for row in tbl_block.rows:
                for c, cell_runs in enumerate(row):
                    if c < n_cols:
                        table.cell(r, c).text = _runs_to_plain_text(cell_runs)
                r += 1

    return prs
