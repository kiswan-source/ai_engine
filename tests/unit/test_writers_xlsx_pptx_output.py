"""write_xlsx/write_pptx + read_xlsx/read_pptx end-to-end (Workspace Slice
3, Fase 12) — same style as test_writers_markdown_output.py: call the
actual public tool functions, not just the shared parser, since a bug could
live in how writers.py wires title/content around the parsed body.
"""
from agent.tools.readers import read_pptx, read_xlsx
from agent.tools.writers import write_pptx, write_xlsx

CONTENT = """# Ringkasan Eksekutif

Dokumen ini **penting** dan berisi catatan produksi.

- Poin satu
- Poin dua

# Data Produksi

| Blok | Luas (Ha) |
| --- | --- |
| Blok A | 11.35 |
| Blok B | 5.20 |
"""


def test_write_xlsx_produces_real_workbook_with_sheets_and_table(tmp_path):
    from openpyxl import load_workbook

    out = tmp_path / "laporan.xlsx"
    result = write_xlsx(str(out), "Laporan Uji", CONTENT)

    assert result["success"] is True
    assert result["type"] == "xlsx"
    wb = load_workbook(str(out))
    assert wb.sheetnames == ["Ringkasan Eksekutif", "Data Produksi"]
    data_ws = wb["Data Produksi"]
    rows = [[c.value for c in row] for row in data_ws.iter_rows()]
    assert rows == [["Blok", "Luas (Ha)"], ["Blok A", "11.35"], ["Blok B", "5.20"]]
    # No literal markdown syntax anywhere in any sheet.
    all_values = [str(c.value) for ws in wb.worksheets for row in ws.iter_rows() for c in row if c.value]
    for leaked in ("**", "| Blok"):
        assert not any(leaked in v for v in all_values)


def test_write_pptx_produces_real_presentation_with_slides_and_table(tmp_path):
    from pptx import Presentation

    out = tmp_path / "laporan.pptx"
    result = write_pptx(str(out), "Laporan Uji", CONTENT)

    assert result["success"] is True
    assert result["type"] == "pptx"
    prs = Presentation(str(out))
    titles = [s.shapes.title.text for s in prs.slides]
    assert titles == ["Ringkasan Eksekutif", "Data Produksi", "Data Produksi — Tabel"]
    table_slide = list(prs.slides)[2]
    table = next(s for s in table_slide.shapes if s.has_table).table
    assert [c.text for c in table.rows[0].cells] == ["Blok", "Luas (Ha)"]


def test_read_xlsx_round_trips_write_xlsx_output(tmp_path):
    out = tmp_path / "laporan.xlsx"
    write_xlsx(str(out), "Laporan Uji", CONTENT)

    result = read_xlsx(str(out))
    assert result["sheet_names"] == ["Ringkasan Eksekutif", "Data Produksi"]
    assert result["sheet_count"] == 2
    assert "Blok A" in result["text"]


def test_read_pptx_round_trips_write_pptx_output(tmp_path):
    out = tmp_path / "laporan.pptx"
    write_pptx(str(out), "Laporan Uji", CONTENT)

    result = read_pptx(str(out))
    assert result["slide_count"] == 3
    assert "Ringkasan Eksekutif" in result["text"]
    assert "Poin satu" in result["text"]


def test_read_xlsx_missing_file_returns_error():
    result = read_xlsx("/no/such/file.xlsx")
    assert "error" in result


def test_read_pptx_missing_file_returns_error():
    result = read_pptx("/no/such/file.pptx")
    assert "error" in result


def test_write_xlsx_reports_failure_not_exception_on_bad_path():
    result = write_xlsx("/no/such/directory/out.xlsx", "T", "isi")
    assert result["success"] is False
    assert "error" in result


def test_write_pptx_reports_failure_not_exception_on_bad_path():
    result = write_pptx("/no/such/directory/out.pptx", "T", "isi")
    assert result["success"] is False
    assert "error" in result
